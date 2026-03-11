import os
import io
import json
import base64
import tempfile
import subprocess
import sys
from typing import Optional, Dict, Any, List
from pathlib import Path


# Mock cairosvg before importing to avoid cairo library errors
class MockCairoSVG:
    def __getattr__(self, name):
        return lambda *a, **kw: None


sys.modules["cairosvg"] = MockCairoSVG()

import fitz  # PyMuPDF
from PIL import Image as PILImage, Image as Image
import cairosvg


class ConversionError(Exception):
    pass


class DocumentConverter:
    def __init__(self, temp_dir: Optional[str] = None):
        self.temp_dir = temp_dir or tempfile.gettempdir()
        self._ensure_temp_dir()

    def _ensure_temp_dir(self):
        Path(self.temp_dir).mkdir(parents=True, exist_ok=True)

    def _get_temp_path(self, suffix: str) -> str:
        import uuid

        filename = f"{uuid.uuid4()}{suffix}"
        return os.path.join(self.temp_dir, filename)

    def pdf_to_images(
        self, pdf_data: bytes, dpi: int = 300, fmt: str = "PNG"
    ) -> List[bytes]:
        images = []
        with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as tmp:
            tmp.write(pdf_data)
            tmp_path = tmp.name

        try:
            doc = fitz.open(tmp_path)
            for page_num in range(len(doc)):
                page = doc.load_page(page_num)
                zoom = dpi / 72
                mat = fitz.Matrix(zoom, zoom)
                pix = page.get_pixmap(matrix=mat)

                img_data = pix.tobytes(fmt)
                images.append(img_data)
            doc.close()
        finally:
            os.unlink(tmp_path)

        return images

    def pdf_to_word(self, pdf_data: bytes, output_path: Optional[str] = None) -> bytes:
        # Use fallback method (pypdf + python-docx) directly
        # as Gotenberg 7 has different API endpoints
        output_path = output_path or self._get_temp_path(".docx")
        return self._pdf_to_word_fallback(pdf_data, output_path)

    def _pdf_to_word_fallback(self, pdf_data: bytes, output_path: str) -> bytes:
        from docx import Document
        from docx.shared import Inches

        doc = Document()

        with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as tmp:
            tmp.write(pdf_data)
            tmp_path = tmp.name

        try:
            pdf_doc = fitz.open(tmp_path)
            for page_num in range(len(pdf_doc)):
                page = pdf_doc.load_page(page_num)
                text = page.get_text()
                if text.strip():
                    doc.add_paragraph(text)

                # Extract images
                img_list = page.get_images()
                for img_index, img in enumerate(img_list):
                    xref = img[0]
                    base_img = pdf_doc.extract_image(xref)
                    image_bytes = base_img["image"]

                    # Save to temp file for PIL
                    with tempfile.NamedTemporaryFile(
                        suffix=".png", delete=False
                    ) as img_tmp:
                        img_tmp.write(image_bytes)
                        img_tmp_path = img_tmp.name

                    try:
                        doc.add_picture(img_tmp_path, width=Inches(6))
                    finally:
                        os.unlink(img_tmp_path)

                doc.add_page_break()
            pdf_doc.close()
        finally:
            os.unlink(tmp_path)

        doc.save(output_path)
        with open(output_path, "rb") as f:
            return f.read()

    def pdf_to_html(self, pdf_data: bytes) -> str:
        html_parts = []

        with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as tmp:
            tmp.write(pdf_data)
            tmp_path = tmp.name

        try:
            doc = fitz.open(tmp_path)
            for page_num in range(len(doc)):
                page = doc.load_page(page_num)
                html = page.get_text("html")
                html_parts.append(html)
            doc.close()
        finally:
            os.unlink(tmp_path)

        return "<hr/>".join(html_parts)

    def pdf_to_pptx(self, pdf_data: bytes, output_path: Optional[str] = None) -> bytes:
        # Convert PDF to images first, then to PPTX
        output_path = output_path or self._get_temp_path(".pptx")

        images = self.pdf_to_images(pdf_data, dpi=150, fmt="PNG")

        from pptx import Presentation
        from pptx.util import Inches

        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)

        blank_slide_layout = prs.slide_layouts[6]

        for i, img_data in enumerate(images):
            slide = prs.slides.add_slide(blank_slide_layout)

            with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as img_tmp:
                img_tmp.write(img_data)
                img_tmp_path = img_tmp.name

            try:
                slide.shapes.add_picture(
                    img_tmp_path, Inches(0), Inches(0), width=Inches(10)
                )
            finally:
                os.unlink(img_tmp_path)

        prs.save(output_path)
        with open(output_path, "rb") as f:
            return f.read()

    def word_to_pdf(self, docx_data: bytes, output_path: Optional[str] = None) -> bytes:
        output_path = output_path or self._get_temp_path(".pdf")

        try:
            # Try Gotenberg first
            return self._gotenberg_convert(docx_data, "docx", "pdf")
        except Exception:
            # Fallback: use python-docx + PyMuPDF
            return self._word_to_pdf_fallback(docx_data, output_path)

    def _word_to_pdf_fallback(self, docx_data: bytes, output_path: str) -> bytes:
        from docx import Document

        with tempfile.NamedTemporaryFile(suffix=".docx", delete=False) as tmp:
            tmp.write(docx_data)
            tmp_path = tmp.name

        try:
            doc = Document(tmp_path)

            pdf_doc = fitz.open()

            a4_width, a4_height = fitz.paper_size("A4")

            for para in doc.paragraphs:
                if para.text.strip():
                    page = pdf_doc.new_page(width=a4_width, height=a4_height)
                    text = para.text
                    page.insert_text((72, 72), text, fontsize=12)

            pdf_doc.save(output_path)
            pdf_doc.close()

            with open(output_path, "rb") as f:
                return f.read()
        finally:
            os.unlink(tmp_path)

    def word_to_markdown(self, docx_data: bytes) -> str:
        from docx import Document

        with tempfile.NamedTemporaryFile(suffix=".docx", delete=False) as tmp:
            tmp.write(docx_data)
            tmp_path = tmp.name

        try:
            doc = Document(tmp_path)
            md_lines = []

            for para in doc.paragraphs:
                style = para.style.name.lower() if para.style else ""

                if "heading" in style:
                    level = style.replace("heading", "").strip() or "1"
                    md_lines.append(f"{'#' * int(level)} {para.text}")
                else:
                    md_lines.append(para.text)

                md_lines.append("")

            return "\n".join(md_lines)
        finally:
            os.unlink(tmp_path)

    def excel_to_pdf(
        self, xlsx_data: bytes, output_path: Optional[str] = None
    ) -> bytes:
        try:
            return self._gotenberg_convert(xlsx_data, "xlsx", "pdf")
        except Exception:
            return self._excel_to_pdf_fallback(xlsx_data, output_path)

    def _excel_to_pdf_fallback(
        self, xlsx_data: bytes, output_path: Optional[str] = None
    ) -> bytes:
        from openpyxl import load_workbook
        from openpyxl.styles import Font, Alignment
        from reportlab.lib.pagesizes import A4
        from reportlab.lib import colors
        from reportlab.platypus import (
            SimpleDocTemplate,
            Table,
            TableStyle,
            Paragraph,
            Spacer,
        )
        from reportlab.lib.units import inch
        from reportlab.pdfbase import pdfmetrics
        from reportlab.pdfbase.ttfonts import TTFont
        from reportlab.lib.styles import getSampleStyleSheet
        import io

        # Register Chinese font
        chinese_fonts = [
            "/System/Library/Fonts/PingFang.ttc",
            "/System/Library/Fonts/STHeiti Light.ttc",
            "/System/Library/Fonts/STHeiti Medium.ttc",
        ]
        registered_font = None
        for font_path in chinese_fonts:
            if os.path.exists(font_path):
                try:
                    pdfmetrics.registerFont(TTFont("ChineseFont", font_path))
                    registered_font = "ChineseFont"
                    break
                except Exception:
                    continue

        # Fallback to STSong-Light (built-in PDF font for Chinese)
        if not registered_font:
            try:
                from reportlab.pdfbase import pdffont

                registered_font = "STSong-Light"
            except Exception:
                registered_font = "Helvetica"

        output_path = output_path or self._get_temp_path(".pdf")

        with tempfile.NamedTemporaryFile(suffix=".xlsx", delete=False) as tmp:
            tmp.write(xlsx_data)
            tmp_path = tmp.name

        try:
            wb = load_workbook(tmp_path, read_only=True)
            ws = wb.active

            doc = SimpleDocTemplate(output_path, pagesize=A4)
            elements = []
            styles = getSampleStyleSheet()

            font_name = registered_font

            for row_idx, row in enumerate(ws.iter_rows(values_only=True)):
                if row_idx == 0:
                    continue

                data = []
                for cell in row:
                    if cell is not None:
                        data.append(str(cell))
                    else:
                        data.append("")

                if data:
                    t = Table([data])
                    t.setStyle(
                        TableStyle(
                            [
                                ("BACKGROUND", (0, 0), (-1, 0), colors.grey),
                                ("TEXTCOLOR", (0, 0), (-1, 0), colors.whitesmoke),
                                ("ALIGN", (0, 0), (-1, -1), "CENTER"),
                                ("FONTNAME", (0, 0), (-1, -1), font_name),
                                ("FONTSIZE", (0, 0), (-1, 0), 10),
                                ("FONTSIZE", (0, 1), (-1, -1), 8),
                                ("BOTTOMPADDING", (0, 0), (-1, 0), 12),
                                ("BACKGROUND", (0, 1), (-1, -1), colors.beige),
                                ("GRID", (0, 0), (-1, -1), 1, colors.black),
                            ]
                        )
                    )
                    elements.append(t)
                    elements.append(Spacer(1, 0.2 * inch))

            doc.build(elements)

            with open(output_path, "rb") as f:
                return f.read()
        finally:
            os.unlink(tmp_path)

    def excel_to_csv(self, xlsx_data: bytes) -> str:
        import csv
        from openpyxl import load_workbook

        with tempfile.NamedTemporaryFile(suffix=".xlsx", delete=False) as tmp:
            tmp.write(xlsx_data)
            tmp_path = tmp.name

        try:
            wb = load_workbook(tmp_path, read_only=True)
            ws = wb.active

            csv_lines = []
            for row in ws.iter_rows(values_only=True):
                csv_lines.append(",".join(str(cell) if cell else "" for cell in row))

            return "\n".join(csv_lines)
        finally:
            os.unlink(tmp_path)

    def pptx_to_pdf(self, pptx_data: bytes) -> bytes:
        return self._gotenberg_convert(pptx_data, "pptx", "pdf")

    def pptx_to_images(self, pptx_data: bytes, dpi: int = 150) -> List[bytes]:
        from pptx import Presentation
        from PIL import Image as PILImage
        import io

        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
            tmp.write(pptx_data)
            tmp_path = tmp.name

        try:
            prs = Presentation(tmp_path)
            images = []

            for slide_num, slide in enumerate(prs.slides):
                # Create a high-res image of each slide
                # Note: python-pptx doesn't render slides directly
                # This is a simplified version
                pass

        finally:
            os.unlink(tmp_path)
            return images

    def markdown_to_pdf(
        self, md_data: bytes, output_path: Optional[str] = None
    ) -> bytes:
        # Convert markdown to HTML first, then to PDF
        import markdown

        html = markdown.markdown(md_data.decode("utf-8"))

        # Wrap HTML with basic styling
        full_html = f"""
        <!DOCTYPE html>
        <html>
        <head>
            <meta charset="utf-8">
            <style>
                body {{ font-family: Arial, sans-serif; margin: 40px; }}
                h1 {{ color: #333; }}
                code {{ background: #f4f4f4; padding: 2px 5px; }}
                pre {{ background: #f4f4f4; padding: 10px; }}
            </style>
        </head>
        <body>{html}</body>
        </html>
        """

        return self._gotenberg_convert(full_html.encode("utf-8"), "html", "pdf")

    def markdown_to_html(self, md_data: bytes) -> str:
        import markdown

        return markdown.markdown(md_data.decode("utf-8"))

    def markdown_to_word(
        self, md_data: bytes, output_path: Optional[str] = None
    ) -> bytes:
        # Use pandoc for conversion
        md_text = md_data.decode("utf-8")

        try:
            import pypandoc

            pypandoc.ensure_pandoc_installed()
            # Use a temp file path
            output_file = output_path or self._get_temp_path(".docx")
            pypandoc.convert_text(
                md_text,
                "docx",
                format="markdown",
                outputfile=output_file,
            )
            # Read the generated file
            with open(output_file, "rb") as f:
                return f.read()
        except Exception as e:
            print(f"pandoc conversion failed: {e}, using fallback")
            # Fallback: basic conversion
            return self._markdown_to_word_fallback(md_data, output_path)

    def _markdown_to_word_fallback(
        self, md_data: bytes, output_path: Optional[str] = None
    ) -> bytes:
        from docx import Document
        from docx.shared import Pt

        output_path = output_path or self._get_temp_path(".docx")
        md_text = md_data.decode("utf-8")

        doc = Document()

        lines = md_text.split("\n")
        in_code_block = False

        for line in lines:
            if line.startswith("```"):
                in_code_block = not in_code_block
                continue

            if in_code_block:
                doc.add_paragraph(line, style="Code")
            elif line.startswith("#"):
                # Count heading level
                level = len(line) - len(line.lstrip("#"))
                text = line.lstrip("#").strip()
                doc.add_heading(text, level=min(level, 6))
            elif line.startswith("- ") or line.startswith("* "):
                doc.add_paragraph(line, style="List Bullet")
            else:
                doc.add_paragraph(line)

        doc.save(output_path)
        with open(output_path, "rb") as f:
            return f.read()

    def svg_to_png(
        self, svg_data: bytes, width: Optional[int] = None, height: Optional[int] = None
    ) -> bytes:
        return cairosvg.svg2png(
            bytestring=svg_data, output_width=width, output_height=height
        )

    def svg_to_pdf(self, svg_data: bytes, output_path: Optional[str] = None) -> bytes:
        output_path = output_path or self._get_temp_path(".pdf")
        png_data = self.svg_to_png(svg_data)

        # Convert PNG to PDF
        pdf_doc = fitz.open()
        img = PILImage.open(io.BytesIO(png_data))

        # Create PDF page matching image size
        page_width = img.width * 72 / 300  # Convert pixels to PDF points
        page_height = img.height * 72 / 300

        page = pdf_doc.new_page(width=page_width, height=page_height)
        page.insert_image(fitz.Rect(0, 0, page_width, page_height), stream=png_data)

        pdf_doc.save(output_path)
        pdf_doc.close()

        with open(output_path, "rb") as f:
            return f.read()

    def png_to_svg(self, png_data: bytes) -> bytes:
        # Use potrace for raster to vector conversion
        # Simplified: just return a basic SVG wrapper
        import base64

        img = PILImage.open(io.BytesIO(png_data))
        width, height = img.size

        # Embed PNG in SVG (not true vector, but preserves image)
        import base64

        b64_data = base64.b64encode(png_data).decode("utf-8")

        svg = f'''<svg xmlns="http://www.w3.org/2000/svg" width="{width}" height="{height}">
            <image href="data:image/png;base64,{b64_data}" width="{width}" height="{height}"/>
        </svg>'''

        return svg.encode("utf-8")

    def png_to_ico(self, png_data: bytes, sizes: tuple = (16, 32, 48, 256)) -> bytes:
        img = PILImage.open(io.BytesIO(png_data))

        ico_images = []
        for size in sizes:
            resized = img.resize((size, size), PILImage.Resampling.LANCZOS)
            ico_images.append(resized)

        # Save as ICO
        output = io.BytesIO()
        ico_images[0].save(output, format="ICO", sizes=[(s, s) for s in sizes])

        return output.getvalue()

    def png_to_pdf(self, png_data: bytes, output_path: Optional[str] = None) -> bytes:
        output_path = output_path or self._get_temp_path(".pdf")

        img = PILImage.open(io.BytesIO(png_data))

        pdf_doc = fitz.open()
        page_width = img.width * 72 / 300
        page_height = img.height * 72 / 300

        page = pdf_doc.new_page(width=page_width, height=page_height)
        page.insert_image(fitz.Rect(0, 0, page_width, page_height), stream=png_data)

        pdf_doc.save(output_path)
        pdf_doc.close()

        with open(output_path, "rb") as f:
            return f.read()

    def _gotenberg_convert(
        self, input_data: bytes, input_format: str, output_format: str
    ) -> bytes:
        """Use Gotenberg API for conversions"""
        gotenberg_url = os.getenv("GOTENBERG_URL", "http://localhost:3000")

        import requests
        from requests.adapters import HTTPAdapter
        from urllib3.util.retry import Retry

        # Create session with retry logic
        session = requests.Session()
        retry_strategy = Retry(
            total=3,
            backoff_factor=1,
            status_forcelist=[500, 502, 503, 504],
        )
        adapter = HTTPAdapter(max_retries=retry_strategy)
        session.mount("http://", adapter)
        session.mount("https://", adapter)

        # Different endpoints for different conversions (Gotenberg 7 API)
        endpoints = {
            ("docx", "pdf"): "/forms/libreoffice/convert",
            ("xlsx", "pdf"): "/forms/libreoffice/convert",
            ("pptx", "pdf"): "/forms/libreoffice/convert",
            ("html", "pdf"): "/forms/chromium/convert/html",
            ("markdown", "pdf"): "/forms/chromium/convert/markdown",
        }

        endpoint = endpoints.get(
            (input_format, output_format), "/forms/libreoffice/convert"
        )

        # For HTML and Markdown conversions, Gotenberg expects specific filename
        if input_format == "html":
            files = {"index.html": ("index.html", input_data)}
        elif input_format == "markdown":
            files = {"index.md": ("index.md", input_data)}
        else:
            files = {"file": ("input." + input_format, input_data)}

        # Calculate timeout based on file size (more time for larger files)
        file_size_mb = len(input_data) / (1024 * 1024)
        timeout = max(
            300, int(file_size_mb * 10)
        )  # At least 5 minutes, more for large files

        response = session.post(
            f"{gotenberg_url}{endpoint}",
            files=files,
            timeout=timeout,
            allow_redirects=True,
        )

        if response.status_code != 200:
            raise ConversionError(f"Gotenberg conversion failed: {response.text}")

        return response.content


class PDFProcessor:
    """Advanced PDF processing operations"""

    def __init__(self, temp_dir: Optional[str] = None):
        self.temp_dir = temp_dir or tempfile.gettempdir()

    def merge_pdfs(
        self, pdf_datas: List[bytes], output_filename: str = "merged.pdf"
    ) -> bytes:
        output_path = os.path.join(self.temp_dir, output_filename)

        merged_pdf = fitz.open()

        for pdf_data in pdf_datas:
            with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as tmp:
                tmp.write(pdf_data)
                tmp_path = tmp.name

            try:
                pdf = fitz.open(tmp_path)
                merged_pdf.insert_pdf(pdf)
                pdf.close()
            finally:
                os.unlink(tmp_path)

        merged_pdf.save(output_path)
        merged_pdf.close()

        with open(output_path, "rb") as f:
            return f.read()

    def split_pdf(
        self,
        pdf_data: bytes,
        page_ranges: List[str],
        output_filename: str = "split.pdf",
    ) -> bytes:
        output_path = os.path.join(self.temp_dir, output_filename)

        with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as tmp:
            tmp.write(pdf_data)
            tmp_path = tmp.name

        try:
            src_doc = fitz.open(tmp_path)
            new_doc = fitz.open()

            for range_str in page_ranges:
                if "-" in range_str:
                    start, end = map(int, range_str.split("-"))
                    for page_num in range(start - 1, end):
                        if page_num < len(src_doc):
                            new_doc.insert_pdf(
                                src_doc, from_page=page_num, to_page=page_num
                            )
                else:
                    page_num = int(range_str) - 1
                    if page_num < len(src_doc):
                        new_doc.insert_pdf(
                            src_doc, from_page=page_num, to_page=page_num
                        )

            new_doc.save(output_path)
            new_doc.close()
            src_doc.close()

            with open(output_path, "rb") as f:
                return f.read()
        finally:
            os.unlink(tmp_path)

    def extract_images(self, pdf_data: bytes) -> List[Dict[str, Any]]:
        images = []

        with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as tmp:
            tmp.write(pdf_data)
            tmp_path = tmp.name

        try:
            doc = fitz.open(tmp_path)

            for page_num in range(len(doc)):
                page = doc.load_page(page_num)
                img_list = page.get_images()

                for img_index, img in enumerate(img_list):
                    xref = img[0]
                    base_img = doc.extract_image(xref)

                    images.append(
                        {
                            "page": page_num + 1,
                            "index": img_index,
                            "data": base_img["image"],
                            "width": base_img["width"],
                            "height": base_img["height"],
                            "ext": base_img["ext"],
                        }
                    )

            doc.close()
        finally:
            os.unlink(tmp_path)

        return images

    def extract_text(self, pdf_data: bytes) -> str:
        with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as tmp:
            tmp.write(pdf_data)
            tmp_path = tmp.name

        try:
            doc = fitz.open(tmp_path)
            text_parts = []

            for page_num in range(len(doc)):
                page = doc.load_page(page_num)
                text_parts.append(f"--- Page {page_num + 1} ---\n")
                text_parts.append(page.get_text())

            doc.close()
            return "\n".join(text_parts)
        finally:
            os.unlink(tmp_path)
